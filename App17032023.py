# impoter les packages
import streamlit as st
import pandas as pd
import numpy as np
from scipy.stats import multivariate_normal
from PIL import Image
from streamlit_option_menu import option_menu
import json
import requests
from streamlit_lottie import st_lottie
import os
import emoji
import time
import base64
from io import StringIO, BytesIO
from pptx import Presentation
from pptx.util import Inches
import sqlite3
import hashlib
# ---------------------- Importer les images -----------------------
image = Image.open('Y3.png')
st.set_page_config(page_title="My Risk Management App",
                   page_icon="üßä",
                   initial_sidebar_state="expanded",
                   layout="wide" # utilise la largeur maximale,
                   )
# Supprimer l'option de menu de streamlit
hide_menu = """
        <style>
        #MainMenu {visibility: hidden;}
        footer {visibility: hidden;}
        </style>
        """
st.markdown(hide_menu, unsafe_allow_html=True)

# definition les fonctions
def make_hashes(password):
    return hashlib.sha256(str.encode(password)).hexdigest()

def check_hashes(password,hashed_text):
    if make_hashes(password) == hashed_text:
        return hashed_text
    return False

# DB Management
conn = sqlite3.connect('data.db')
c = conn.cursor()
# DB  Functions
def create_usertable():
    c.execute('CREATE TABLE IF NOT EXISTS userstable(username TEXT,password TEXT)')


def add_userdata(username,password):
    c.execute('INSERT INTO userstable(username,password) VALUES (?,?)',(username,password))
    conn.commit()

def login_user(username,password):
    c.execute('SELECT * FROM userstable WHERE username =? AND password = ?',(username,password))
    data = c.fetchall()
    return data


def view_all_users():
    c.execute('SELECT * FROM userstable')
    data = c.fetchall()
    return data


def main():
    # Cr√©er un expander dans la barre lat√©rale
    options_expander = st.sidebar.expander("Je me connecte")
    with options_expander:
        choice=st.selectbox("",["S'inscrire","Se connecter"],index=1)
    if choice=="Se connecter":
        with options_expander:
            username = st.text_input("Identifiant")
            password = st.text_input("Mot de passe",type='password')
        if username and password is not None:
            if st.sidebar.checkbox("Se connecter"):
                create_usertable()
                hashed_pswd = make_hashes(password)

                result = login_user(username,check_hashes(password,hashed_pswd))
                if result:
                    with st.spinner("Conexion en cours..."):
                        time.sleep(1)
                        result = st.success("Connxion effectu√©e avec succ√®s")
                        # Attendre 3 secondes avant de masquer l'alerte
                        for i in range(1):
                            time.sleep(1)

                        # Effacer l'alerte de succ√®s
                        result.empty()

                    # ---------------------Option de d√©connexion
                    if st.button("D√©connexion"):
                        st.success("D√©connect√© avec succ√®s !.")
                        st.stop()

                    test_a_faire= option_menu(
                        menu_title=None,
                        options=["Home","Stress Tests","Reverses Tests","IFRS 9","Contact"],
                        icons=["house","list-task","list-task","list-task",'envelope'],
                        menu_icon="cast",
                        default_index=0,
                        orientation="horizontal",

                        styles={
                            "container": {"padding": "0!important","background-color": "#cccccc"},
                            "icon": {"color": "white", "font-size": "16px"},
                            "nav-link": {"font-size": "16px", "text-align": "center", "margin":"0px"},
                            "nav-link-selected": {"background-color": "#0E6655"},
                        }
                    )

                    if test_a_faire=="Home":
                        st.cache(allow_output_mutation=True,show_spinner=True,suppress_st_warning=True)
                        def load_lottiefile(filepath: str):
                            with open(filepath, "r") as f:
                                return json.load(f)
                        def load_lottieurl(url: str):
                            r = requests.get(url)
                            if r.status_code != 200:
                                return None
                            return r.json()
                        z1,sep,z2=st.columns([8,0.2,8])
                        st.markdown("***")
                        with z1:
                            st.markdown("""
                                  ***
                                 ### Welcome to the Risk Management app
                                  
                                  ***
                                  """)

                        with z2:
                            st.image(image, caption="Audits & Conseils")
                            st.write(emoji.emojize("""# Risk Management App """))
                        lottie_hello = load_lottiefile("110031-welcome.json")
                        # lottie_hello = load_lottieurl("https://assets1.lottiefiles.com/packages/lf20_d8cmv6qa.json")
                        with z1:
                            st_lottie(
                                lottie_hello,
                                speed=1,
                                reverse=False,
                                loop=True,
                                quality="low", # medium ; high
                                #renderer="svg", # canvas
                                height=None,
                                width=None,
                                key=None,
                            )
                            st.balloons()


                    elif test_a_faire == "Contact":

                        st.cache(allow_output_mutation=True,show_spinner=True,suppress_st_warning=True)
                        e1, sep,e2=st.columns([8,0.2,8])
                        with e1:
                            """
                            ## Departement Risk Management
                            ### Etude Statistique et mod√©lisation
                            ***
                            #### Sylvain N'DRI
                            * Economiste | Econom√®tre Financier | Data Analyst & Machine Learning Consultant Junior
                            * Contact :+225 0778804983/0140015896
                            * Email: sylvain.ngoran.ndri@gmail.com
                            * Vous pouvez acc√©der √† mon [linkedin](https://www.linkedin.com/in/n-goran-sylvain-n-dri-1b778b1a1/) en cliquant sur ce lien.
                            ***
                            """
                        with e2 :
                            def load_lottiefile(filepath: str):
                                with open(filepath, "r") as f:
                                    return json.load(f)
                            def load_lottieurl(url: str):
                                r = requests.get(url)
                                if r.status_code != 200:
                                    return None
                                return r.json()

                            lottie_hello = load_lottiefile("h.json")
                            #lottie_hello = load_lottieurl("https://assets1.lottiefiles.com/private_files/lf30_TBKozE.json")
                            # st.markdown("***")
                            st_lottie(
                                lottie_hello,
                                speed=1,
                                reverse=False,
                                loop=True,
                                quality="low", # medium ; high
                                #renderer=
                                height=None,
                                width=None,
                                key=None,
                            )



                    elif test_a_faire== "Stress Tests":
                        global  df
                        st.cache(allow_output_mutation=True,show_spinner=True,suppress_st_warning=True)
                        st.title("STRESS TEST : RISK MANAGEMENT TOOL ")
                        st.markdown("***")

                        with open('style.css') as f:
                            st.markdown(f'<style>{f.read()}</style>', unsafe_allow_html=True)

                        #  -------------------------------------------------------------

                        # model 1 pnp
                        model1_coef_1,model1_coef_2,model1_coef_3,model1_coef_4,model1_coef_5,model1_coef_6=[2,3,4,1.5,4,4]
                        cons1=15.00

                        # prevision function
                        # MODEL 1 ========================================================================================================

                        sce_1, sep, sce_2 ,sep,sce_3,sep,sce_4= st.columns([6,0.2,6,0.2,6,0.2,6])
                        with sce_2:
                            scenario1=st.number_input("Variation 1: excemple PIB", 0.05)
                        with sce_3:
                            scenario2=st.number_input("Variation 2: exemple inflation", 0.05)

                        options_expander=st.sidebar.expander("Je param√®tre le mod√®le 1")
                        with options_expander:
                            st.header("Mod√®le 1")
                            sc_1, sep, sc_2 ,sep,sc_3,= st.columns([4,0.2,4,0.2,4])
                            with sc_1:
                                m1_variable_1=st.slider("Variable 1", 1.00,30.00)/100
                                m1_variable_2=st.slider("Variable 2", 0.01,50.00)/100
                                m1_variable_3=st.slider("Variable 3", 1.00,100.00)/100
                            with sc_2:
                                m1_variable_4=st.slider("Variable 4", 1.00,100.00)/100
                                m1_variable_5=st.slider("Variable 5", 1.00,100.00)/100
                                m1_variable_6=st.slider("Variable 6", 1.00,100.00)/100

                        # MODEL 2 =================================================================================================================
                        # model 2 cex
                        model2_coef_1,model2_coef_2,model2_coef_3,model2_coef_4,model2_coef_5,model2_coef_6=[2,3,4,1.5,4,4]
                        cons2=70.00

                        options_expander=st.sidebar.expander("Je param√®tre le mod√®le 2")
                        with options_expander:
                            sc_1, sep, sc_2 ,sep,sc_3,= st.columns([4,0.2,4,0.2,4])
                            with sc_1:
                                st.header("Mod√®le 2")
                                m2_variable_1=st.slider("M2 Variable 1", 1.00,50.00)/100
                                m2_variable_2=st.slider("M2 Variable 2", 1.00,50.00)/100
                                m2_variable_3=st.slider("M2 Variable 3", 1.00,50.00)/100
                            with sc_2:
                                st.markdown("""
                                            #
                                            #
                                            """)
                                m2_variable_4=st.slider("M2 Variable 4", 1.00,50.00)/100
                                m2_variable_5=st.slider("M2 Variable 5", 1.00,50.00)/100
                                m2_variable_6=st.slider("M2 Variable 6", 1.00,50.00)/100

                        # MODEL 3=============================================================================================
                        #===================================================================
                        model3_coef_1,model3_coef_2,model3_coef_3,model3_coef_4,model3_coef_5,model3_coef_6=[2,3,4,1.5,4,4]
                        cons3=13.49
                        options_expander=st.sidebar.expander("Je param√®tre le mod√®le 3")
                        with options_expander:
                            sc_1, sep, sc_2 ,sep,sc_3,= st.columns([4,0.2,4,0.2,4])
                            with sc_1:
                                st.header("Mod√®le 3")
                                m3_variable_1=st.slider("M3 Variable 1", 1.00,50.00)/100
                                m3_variable_2=st.slider("M3 Variable 2", 1.00,50.00)/100
                                m3_variable_3=st.slider("M3 Variable 3", 1.00,50.00)/100
                            with sc_2:
                                st.markdown("""
                                            #
                                            #
                                            """)
                                m3_variable_4=st.slider("M3 Variable 4", 1.00,100.00)/100
                                m3_variable_5=st.slider("M3 Variable 5", 1.00,50.00)/100
                                m3_variable_6=st.slider("M3 Variable 6", 1.00,50.00)/100

                        #  """
                        #  Scenario 3
                        #   """


                        # Scenario 3 model 1 ===================================================================
                        data={"var1":m1_variable_1+scenario1,"var2":m1_variable_2+scenario2,"var3":m1_variable_3,"var4":m1_variable_4,
                              "var5":m1_variable_5,"var6":m1_variable_6}
                        datacof={"c1":model1_coef_1,"c2":model1_coef_2,"c3":model1_coef_3,"c4":model1_coef_4,
                                 "c5":model1_coef_5,"c6":model1_coef_6}
                        data=data.values()
                        datacof=datacof.values()
                        data_v0=sum([i*j for (i, j) in zip(data, datacof)])
                        data_v0=data_v0+cons1

                        data_1={"var1":m1_variable_1+scenario1*2,"var2":m1_variable_2+scenario2*2,"var3":m1_variable_3,"var4":m1_variable_4,
                                "var5":m1_variable_5,"var6":m1_variable_6}
                        data_v_1=data_1.values()
                        data_v1=sum([i*j for (i, j) in zip(data_v_1, datacof)])
                        data_v1=data_v1+cons1

                        data_2={"var1":m1_variable_1+scenario1*3,"var2":m1_variable_2+scenario2*3,"var3":m1_variable_3,"var4":m1_variable_4,
                                "var5":m1_variable_5,"var6":m1_variable_6}
                        data_v_2=data_2.values()
                        data_v2=sum([i*j for (i, j) in zip(data_v_2, datacof)])
                        data_v2=data_v2+cons1

                        data_3={"var1":m1_variable_1+scenario1*4,"var2":m1_variable_2+scenario2*4,"var3":m1_variable_3,"var4":m1_variable_4,
                                "var5":m1_variable_5,"var6":m1_variable_6}
                        data_v_3=data_3.values()
                        data_v3=sum([i*j for (i, j) in zip(data_v_3, datacof)])
                        data_v3=data_v3+cons1

                        list_predict_value=[data_v0,data_v1,data_v2,data_v3]
                        dataframe=pd.DataFrame(["Ann√©e 1","Ann√©e 2","Ann√©e 3","Ann√©e 4"],columns=["Ann√©e"])

                        # scenario 3 model 2==============================================================================================

                        data_m2_v0={"var1":m2_variable_1+scenario1,"var2":m2_variable_2+scenario2,"var3":m2_variable_3,"var4":m2_variable_4,
                                    "var5":m2_variable_5,"var6":m2_variable_6}
                        datacof_m2={"c1":model2_coef_1,"c2":model2_coef_2,"c3":model2_coef_3,"c4":model2_coef_4,
                                    "c5":model2_coef_5,"c6":model2_coef_6}
                        data_m2_v0=data_m2_v0.values()
                        datacof_m2=datacof_m2.values()
                        data_m2_v0=sum([i*j for (i, j) in zip(data_m2_v0, datacof_m2)])
                        data_m2_v0=data_m2_v0+cons2
                        data_m2_v1={"var1":m2_variable_1+scenario1*2,"var2":m2_variable_2+scenario2*2,"var3":m2_variable_3,"var4":m2_variable_4,
                                    "var5":m2_variable_5,"var6":m2_variable_6}
                        data_m2_v1=data_m2_v1.values()
                        data_m2_v1=sum([i*j for (i, j) in zip(data_m2_v1, datacof_m2)])
                        data_m2_v1=data_m2_v1+cons2

                        data_m2_v2={"var1":m2_variable_1+scenario1*3,"var2":m2_variable_2+scenario2*3,"var3":m2_variable_3,"var4":m2_variable_4,
                                    "var5":m2_variable_5,"var6":m2_variable_6}
                        data_m2_v2=data_m2_v2.values()
                        data_m2_v2=sum([i*j for (i, j) in zip(data_m2_v2, datacof_m2)])
                        data_m2_v2=data_m2_v2+cons2

                        data_m2_v3={"var1":m2_variable_1+scenario1*4,"var2":m2_variable_2+scenario2*4,"var3":m2_variable_3,"var4":m2_variable_4,
                                    "var5":m2_variable_5,"var6":m2_variable_6}
                        data_m2_v3=data_m2_v3.values()
                        data_m2_v3=sum([i*j for (i, j) in zip(data_m2_v3, datacof_m2)])
                        data_m2_v3=data_m2_v3+cons2
                        list_predict_m2=[data_m2_v0,data_m2_v1,data_m2_v2,data_m2_v3]

                        # scenario 3 model 3 =========================================================

                        data_m3_v0={"var1":m3_variable_1+scenario1,"var2":m3_variable_2+scenario2,"var3":m3_variable_3,"var4":m3_variable_4,
                                    "var5":m3_variable_5,"var6":m3_variable_6}
                        datacof_m3={"c1":model3_coef_1,"c2":model3_coef_2,"c3":model3_coef_3,"c4":model3_coef_4,
                                    "c5":model3_coef_5,"c6":model3_coef_6}
                        data_m3_v0=data_m3_v0.values()
                        datacof_m3=datacof_m3.values()
                        data_m3_v0=sum([i*j for (i, j) in zip(data_m3_v0, datacof_m3)])
                        data_m3_v0=data_m3_v0+cons3

                        data_m3_v1={"var1":m3_variable_1+scenario1*2,"var2":m3_variable_2+scenario2*2,"var3":m3_variable_3,"var4":m3_variable_4,
                                    "var5":m3_variable_5,"var6":m3_variable_6}
                        data_m3_v1=data_m3_v1.values()
                        data_m3_v1=sum([i*j for (i, j) in zip(data_m3_v1, datacof_m3)])
                        data_m3_v1=data_m3_v1+cons3

                        data_m3_v2={"var1":m3_variable_1+scenario1*3,"var2":m3_variable_2+scenario2*3,"var3":m3_variable_3,"var4":m3_variable_4,
                                    "var5":m3_variable_5,"var6":m3_variable_6}
                        data_m3_v2=data_m3_v2.values()
                        data_m3_v2=sum([i*j for (i, j) in zip(data_m3_v2, datacof_m3)])
                        data_m3_v2=data_m3_v2+cons3

                        data_m3_v3={"var1":m3_variable_1+scenario1*4,"var2":m3_variable_2+scenario2*4,"var3":m3_variable_3,"var4":m3_variable_4,
                                    "var5":m3_variable_5,"var6":m3_variable_6}
                        data_m3_v3=data_m3_v3.values()
                        data_m3_v3=sum([i*j for (i, j) in zip(data_m3_v3, datacof_m3)])
                        data_m3_v3=data_m3_v3+cons3
                        list_predict_m3=[data_m3_v0,data_m3_v1,data_m3_v2,data_m3_v3]
                        ##"""

                        #### Scenario 2

                        ##"""
                        # scenario 2 model 1
                        data_s2={"var1":m1_variable_1,"var2":m1_variable_2+scenario2,"var3":m1_variable_3,"var4":m1_variable_4,
                                 "var5":m1_variable_5,"var6":m1_variable_6}
                        datacof_s2={"c1":model1_coef_1,"c2":model1_coef_2,"c3":model1_coef_3,"c4":model1_coef_4,
                                    "c5":model1_coef_5,"c6":model1_coef_6}
                        data_s2=data_s2.values()
                        datacof_s2=datacof_s2.values()
                        data_v0_s2=sum([i*j for (i, j) in zip(data_s2, datacof_s2)])
                        data_v0_s2=data_v0_s2+cons1

                        data_1_s2={"var1":m1_variable_1,"var2":m1_variable_2+scenario2*2,"var3":m1_variable_3,"var4":m1_variable_4,
                                   "var5":m1_variable_5,"var6":m1_variable_6}
                        data_v_1_s2=data_1_s2.values()
                        data_v1_s2=sum([i*j for (i, j) in zip(data_v_1_s2, datacof_s2)])
                        data_v1_s2=data_v1_s2+cons1

                        data_2_s2={"var1":m1_variable_1,"var2":m1_variable_2+scenario2*3,"var3":m1_variable_3,"var4":m1_variable_4,
                                   "var5":m1_variable_5,"var6":m1_variable_6}
                        data_v_2_s2=data_2_s2.values()
                        data_v2_s2=sum([i*j for (i, j) in zip(data_v_2_s2, datacof_s2)])
                        data_v2_s2=data_v2_s2+cons1

                        data_3_s2={"var1":m1_variable_1,"var2":m1_variable_2+scenario2*4,"var3":m1_variable_3,"var4":m1_variable_4,
                                   "var5":m1_variable_5,"var6":m1_variable_6}
                        data_v_3_s2=data_3_s2.values()
                        data_v3_s2=sum([i*j for (i, j) in zip(data_v_3_s2, datacof_s2)])
                        data_v3_s2=data_v3_s2+cons1

                        list_predict_value_s2=[data_v0_s2,data_v1_s2,data_v2_s2,data_v3_s2]


                        # scenario 2 Model 2 ============================================================================================================

                        data_m2_v0_s2={"var1":m2_variable_1,"var2":m2_variable_2+scenario2,"var3":m2_variable_3,"var4":m2_variable_4,
                                       "var5":m2_variable_5,"var6":m2_variable_6}
                        datacof_m2_s2={"c1":model2_coef_1,"c2":model2_coef_2,"c3":model2_coef_3,"c4":model2_coef_4,
                                       "c5":model2_coef_5,"c6":model2_coef_6}
                        data_m2_v0_s2=data_m2_v0_s2.values()
                        datacof_m2_s2=datacof_m2_s2.values()
                        data_m2_v0_s2=sum([i*j for (i, j) in zip(data_m2_v0_s2, datacof_m2_s2)])
                        data_m2_v0_s2=data_m2_v0_s2+cons2

                        data_m2_v1_s2={"var1":m2_variable_1,"var2":m2_variable_2+scenario2*2,"var3":m2_variable_3,"var4":m2_variable_4,
                                       "var5":m2_variable_5,"var6":m2_variable_6}
                        data_m2_v1_s2=data_m2_v1_s2.values()
                        data_m2_v1_s2=sum([i*j for (i, j) in zip(data_m2_v1_s2, datacof_m2_s2)])
                        data_m2_v1_s2=data_m2_v1_s2+cons2

                        data_m2_v2_s2={"var1":m2_variable_1,"var2":m2_variable_2+scenario2*3,"var3":m2_variable_3,"var4":m2_variable_4,
                                       "var5":m2_variable_5,"var6":m2_variable_6}
                        data_m2_v2_s2=data_m2_v2_s2.values()
                        data_m2_v2_s2=sum([i*j for (i, j) in zip(data_m2_v2_s2, datacof_m2_s2)])
                        data_m2_v2_s2=data_m2_v2_s2+cons2

                        data_m2_v3_s2={"var1":m2_variable_1,"var2":m2_variable_2+scenario2*4,"var3":m2_variable_3,"var4":m2_variable_4,
                                       "var5":m2_variable_5,"var6":m2_variable_6}
                        data_m2_v3_s2=data_m2_v3_s2.values()
                        data_m2_v3_s2=sum([i*j for (i, j) in zip(data_m2_v3_s2, datacof_m2_s2)])
                        data_m2_v3_s2=data_m2_v3_s2+cons2
                        list_predict_m2_s2=[data_m2_v0_s2,data_m2_v1_s2,data_m2_v2_s2,data_m2_v3_s2]

                        # scenario 2 model 3

                        data_m3_v0_s2={"var1":m3_variable_1,"var2":m3_variable_2+scenario2,"var3":m3_variable_3,"var4":m3_variable_4,
                                       "var5":m3_variable_5,"var6":m3_variable_6}
                        datacof_m3_s2={"c1":model3_coef_1,"c2":model3_coef_2,"c3":model3_coef_3,"c4":model3_coef_4,
                                       "c5":model3_coef_5,"c6":model3_coef_6}
                        data_m3_v0_s2=data_m3_v0_s2.values()
                        datacof_m3_s2=datacof_m3_s2.values()
                        data_m3_v0_s2=sum([i*j for (i, j) in zip(data_m3_v0_s2, datacof_m3_s2)])
                        data_m3_v0_s2=data_m3_v0_s2+cons3

                        data_m3_v1_s2={"var1":m3_variable_1,"var2":m3_variable_2+scenario2*2,"var3":m3_variable_3,"var4":m3_variable_4,
                                       "var5":m3_variable_5,"var6":m3_variable_6}
                        data_m3_v1_s2=data_m3_v1_s2.values()
                        data_m3_v1_s2=sum([i*j for (i, j) in zip(data_m3_v1_s2, datacof_m3_s2)])
                        data_m3_v1_s2=data_m3_v1_s2+cons3

                        data_m3_v2_s2={"var1":m3_variable_1,"var2":m3_variable_2+scenario2*3,"var3":m3_variable_3,"var4":m3_variable_4,
                                       "var5":m3_variable_5,"var6":m3_variable_6}
                        data_m3_v2_s2=data_m3_v2_s2.values()
                        data_m3_v2_s2=sum([i*j for (i, j) in zip(data_m3_v2_s2, datacof_m3_s2)])
                        data_m3_v2_s2=data_m3_v2_s2+cons3

                        data_m3_v3_s2={"var1":m3_variable_1,"var2":m3_variable_2+scenario2*4,"var3":m3_variable_3,"var4":m3_variable_4,
                                       "var5":m3_variable_5,"var6":m3_variable_6}
                        data_m3_v3_s2=data_m3_v3_s2.values()
                        data_m3_v3_s2=sum([i*j for (i, j) in zip(data_m3_v3_s2, datacof_m3_s2)])
                        data_m3_v3_s2=data_m3_v3_s2+cons3
                        list_predict_m3_s2=[data_m3_v0_s2,data_m3_v1_s2,data_m3_v2_s2,data_m3_v3_s2]


                        ## Sc√©nario 1================
                        #
                        #=============================
                        # scenario 1 model 1
                        data_s1={"var1":m1_variable_1+scenario1,"var2":m1_variable_2,"var3":m1_variable_3,"var4":m1_variable_4,
                                 "var5":m1_variable_5,"var6":m1_variable_6}
                        datacof_s1={"c1":model1_coef_1,"c2":model1_coef_2,"c3":model1_coef_3,"c4":model1_coef_4,
                                    "c5":model1_coef_5,"c6":model1_coef_6}
                        data_s1=data_s1.values()
                        datacof_s1=datacof_s1.values()
                        data_v0_s1=sum([i*j for (i, j) in zip(data_s1, datacof_s1)])
                        data_v0_s1=data_v0_s1+cons1

                        data_1_s1={"var1":m1_variable_1+scenario1*2,"var2":m1_variable_2,"var3":m1_variable_3,"var4":m1_variable_4,
                                   "var5":m1_variable_5,"var6":m1_variable_6}
                        data_v_1_s1=data_1_s1.values()
                        data_v1_s1=sum([i*j for (i, j) in zip(data_v_1_s1, datacof_s1)])
                        data_v1_s1=data_v1_s1+cons1

                        data_2_s1={"var1":m1_variable_1+scenario1*3,"var2":m1_variable_2,"var3":m1_variable_3,"var4":m1_variable_4,
                                   "var5":m1_variable_5,"var6":m1_variable_6}
                        data_v_2_s1=data_2_s1.values()
                        data_v2_s1=sum([i*j for (i, j) in zip(data_v_2_s1, datacof_s1)])
                        data_v2_s1=data_v2_s1+cons1

                        data_3_s1={"var1":m1_variable_1+scenario1*4,"var2":m1_variable_2,"var3":m1_variable_3,"var4":m1_variable_4,
                                   "var5":m1_variable_5,"var6":m1_variable_6}
                        data_v_3_s1=data_3_s1.values()
                        data_v3_s1=sum([i*j for (i, j) in zip(data_v_3_s1, datacof_s1)])
                        data_v3_s1=data_v3_s1+cons1
                        list_predict_value_s1=[data_v0_s1,data_v1_s1,data_v2_s1,data_v3_s1]

                        # scenario 1 model 2

                        data_m2_v0_s1={"var1":m2_variable_1+scenario1,"var2":m2_variable_2,"var3":m2_variable_3,"var4":m2_variable_4,
                                       "var5":m2_variable_5,"var6":m2_variable_6}
                        datacof_m2_s1={"c1":model2_coef_1,"c2":model2_coef_2,"c3":model2_coef_3,"c4":model2_coef_4,
                                       "c5":model2_coef_5,"c6":model2_coef_6}
                        data_m2_v0_s1=data_m2_v0_s1.values()
                        datacof_m2_s1=datacof_m2_s1.values()
                        data_m2_v0_s1=sum([i*j for (i, j) in zip(data_m2_v0_s1, datacof_m2_s1)])
                        data_m2_v0_s1=data_m2_v0_s1+cons2

                        data_m2_v1_s1={"var1":m2_variable_1+scenario1*2,"var2":m2_variable_2,"var3":m2_variable_3,"var4":m2_variable_4,
                                       "var5":m2_variable_5,"var6":m2_variable_6}
                        data_m2_v1_s1=data_m2_v1_s1.values()
                        data_m2_v1_s1=sum([i*j for (i, j) in zip(data_m2_v1_s1, datacof_m2_s1)])
                        data_m2_v1_s1=data_m2_v1_s1+cons2

                        data_m2_v2_s1={"var1":m2_variable_1+scenario1*3,"var2":m2_variable_2,"var3":m2_variable_3,"var4":m2_variable_4,
                                       "var5":m2_variable_5,"var6":m2_variable_6}
                        data_m2_v2_s1=data_m2_v2_s1.values()
                        data_m2_v2_s1=sum([i*j for (i, j) in zip(data_m2_v2_s1, datacof_m2_s1)])
                        data_m2_v2_s1=data_m2_v2_s1+cons2

                        data_m2_v3_s1={"var1":m2_variable_1+scenario1*4,"var2":m2_variable_2,"var3":m2_variable_3,"var4":m2_variable_4,
                                       "var5":m2_variable_5,"var6":m2_variable_6}
                        data_m2_v3_s1=data_m2_v3_s1.values()
                        data_m2_v3_s1=sum([i*j for (i, j) in zip(data_m2_v3_s1, datacof_m2_s1)])
                        data_m2_v3_s1=data_m2_v3_s1+cons2
                        list_predict_m2_s1=[data_m2_v0_s1,data_m2_v1_s1,data_m2_v2_s1,data_m2_v3_s1]

                        #Scenario 1 model 3


                        data_m3_v0_s1={"var1":m3_variable_1+scenario1,"var2":m3_variable_2,"var3":m3_variable_3,"var4":m3_variable_4,
                                       "var5":m3_variable_5,"var6":m3_variable_6}
                        datacof_m3_s1={"c1":model3_coef_1,"c2":model3_coef_2,"c3":model3_coef_3,"c4":model3_coef_4,
                                       "c5":model3_coef_5,"c6":model3_coef_6}
                        data_m3_v0_s1=data_m3_v0_s1.values()
                        datacof_m3_s1=datacof_m3_s1.values()
                        data_m3_v0_s1=sum([i*j for (i, j) in zip(data_m3_v0_s1, datacof_m3_s1)])
                        data_m3_v0_s1=data_m3_v0_s1+cons3

                        data_m3_v1_s1={"var1":m3_variable_1+scenario1*2,"var2":m3_variable_2,"var3":m3_variable_3,"var4":m3_variable_4,
                                       "var5":m3_variable_5,"var6":m3_variable_6}
                        data_m3_v1_s1=data_m3_v1_s1.values()
                        data_m3_v1_s1=sum([i*j for (i, j) in zip(data_m3_v1_s1, datacof_m3_s1)])
                        data_m3_v1_s1=data_m3_v1_s1+cons3

                        data_m3_v2_s1={"var1":m3_variable_1+scenario1*3,"var2":m3_variable_2,"var3":m3_variable_3,"var4":m3_variable_4,
                                       "var5":m3_variable_5,"var6":m3_variable_6}
                        data_m3_v2_s1=data_m3_v2_s1.values()
                        data_m3_v2_s1=sum([i*j for (i, j) in zip(data_m3_v2_s1, datacof_m3_s1)])
                        data_m3_v2_s1=data_m3_v2_s1+cons3

                        data_m3_v3_s1={"var1":m3_variable_1+scenario1*4,"var2":m3_variable_2,"var3":m3_variable_3,"var4":m3_variable_4,
                                       "var5":m3_variable_5,"var6":m3_variable_6}
                        data_m3_v3_s1=data_m3_v3_s1.values()
                        data_m3_v3_s1=sum([i*j for (i, j) in zip(data_m3_v3_s1, datacof_m3_s1)])
                        data_m3_v3_s1=data_m3_v3_s1+cons3
                        list_predict_m3_s1=[data_m3_v0_s1,data_m3_v1_s1,data_m3_v2_s1,data_m3_v3_s1]


                        ### """

                        # ETAPE DE CALCUL SCENARIO 3

                        ###"""


                        df=pd.DataFrame(list_predict_value,index=None,columns=["PNP"])
                        dataframe["PNP 0"]=9.37
                        dataframe["CEX 0"]=69.03
                        dataframe["TCH 0"]=13.88
                        dataframe["Valeur PNP 0"]=910200
                        dataframe["Valeur CEX 0"]=639206
                        dataframe["Valeur TCH 0"]=554.53
                        dataframe["PNP"]=df["PNP"]
                        dfm2=pd.DataFrame(list_predict_m2,index=None, columns=["CEX"])
                        dataframe["CEX"]=dfm2["CEX"]
                        dfm3=pd.DataFrame(list_predict_m3,index=None,columns=["TCH"])
                        dataframe["TCH"]=dfm3["TCH"]

                        # variation
                        def variation(value_one,value_two):
                            return (value_one-value_two)/value_two
                        dataframe["Variation PNP"]=dataframe[["PNP","PNP 0"]].apply(lambda dataframe:variation(dataframe["PNP"],dataframe["PNP 0"]),axis=1)
                        dataframe["Variation CEX"]=dataframe[["CEX","CEX 0"]].apply(lambda dataframe:variation(dataframe["CEX"],dataframe["CEX 0"]),axis=1)
                        dataframe["Variation TCH"]=dataframe[["TCH","TCH 0"]].apply(lambda dataframe:variation(dataframe["TCH"],dataframe["TCH 0"]),axis=1)
                        # prediction en valeur/ value predict
                        def actualisate(value_one,value_two):
                            return value_one*(1+value_two)
                        dataframe["Valeur PNP"]=dataframe[["Valeur PNP 0","Variation PNP"]].apply(lambda dataframe:actualisate(dataframe["Valeur PNP 0"],
                                                                                                                               dataframe["Variation PNP"]),axis=1)
                        dataframe["Valeur CEX"]=dataframe[["Valeur CEX 0","Variation CEX"]].apply(lambda dataframe:actualisate(dataframe["Valeur CEX 0"],
                                                                                                                               dataframe["Variation CEX"]),axis=1)
                        dataframe["Valeur TCH"]=dataframe[["Valeur TCH 0","Variation TCH"]].apply(lambda dataframe:actualisate(dataframe["Valeur TCH 0"],
                                                                                                                               dataframe["Variation TCH"]),axis=1)

                        # risque supplementaire/ sup risk
                        def risk(value_one, value_two):
                            return  value_one - value_two
                        dataframe["risque sup PNP"]=dataframe[["Valeur PNP","Valeur PNP 0"]].apply(lambda dataframe: risk(dataframe["Valeur PNP"],
                                                                                                                          dataframe["Valeur PNP 0"]),axis=1)
                        dataframe["risque sup CEX"]=dataframe[["Valeur CEX","Valeur CEX 0"]].apply(lambda dataframe: risk(dataframe["Valeur CEX"],
                                                                                                                          dataframe["Valeur CEX 0"]),axis=1)

                        dataframe["risque sup TCH"]=dataframe[["Valeur TCH","Valeur TCH 0"]].apply(lambda dataframe: risk(dataframe["Valeur TCH"],
                                                                                                                          dataframe["Valeur TCH 0"]),axis=1)

                        # Calcul des risques  ponderes
                        dataframe["RWA 0"]=10063541
                        dataframe["RWA Pr√©dicte"]=dataframe["risque sup PNP"]*0.85+dataframe["risque sup CEX"]*0.12+dataframe["risque sup TCH"]*0.03+ dataframe["RWA 0"]
                        # ratio de solvabilit√©
                        # st.write("Sc√©nario 3",dataframe)

                        ###"""
                        #### METHODE DE CALCUL SCENARIO 2

                        ####

                        df_s2=pd.DataFrame(list_predict_value_s2,index=None,columns=["PNP S2"])
                        # st.write(df_s2)
                        dataframe["PNP S2"]=df_s2["PNP S2"]
                        dfm2_s2=pd.DataFrame(list_predict_m2_s2,index=None, columns=["CEX S2"])
                        dataframe["CEX S2"]=dfm2_s2["CEX S2"]
                        dfm3_s2=pd.DataFrame(list_predict_m3_s2,index=None,columns=["TCH S2"])
                        dataframe["TCH S2"]=dfm3_s2["TCH S2"]

                        dataframe["Variation PNP S2"]=dataframe[["PNP S2","PNP 0"]].apply(lambda dataframe:variation(dataframe["PNP S2"],dataframe["PNP 0"]),axis=1)
                        dataframe["Variation CEX S2"]=dataframe[["CEX S2","CEX 0"]].apply(lambda dataframe:variation(dataframe["CEX S2"],dataframe["CEX 0"]),axis=1)
                        dataframe["Variation TCH S2"]=dataframe[["TCH S2","TCH 0"]].apply(lambda dataframe:variation(dataframe["TCH S2"],dataframe["TCH 0"]),axis=1)

                        # prediction en valeur/ value predict

                        dataframe["Valeur PNP S2"]=dataframe[["Valeur PNP 0","Variation PNP S2"]].apply(lambda dataframe:actualisate(dataframe["Valeur PNP 0"],
                                                                                                                                     dataframe["Variation PNP S2"]),axis=1)
                        dataframe["Valeur CEX S2"]=dataframe[["Valeur CEX 0","Variation CEX S2"]].apply(lambda dataframe:actualisate(dataframe["Valeur CEX 0"],
                                                                                                                                     dataframe["Variation CEX S2"]),axis=1)
                        dataframe["Valeur TCH S2"]=dataframe[["Valeur TCH 0","Variation TCH S2"]].apply(lambda dataframe:actualisate(dataframe["Valeur TCH 0"],
                                                                                                                                     dataframe["Variation TCH S2"]),axis=1)

                        # risque supplementaire/ sup risk

                        dataframe["risque sup PNP S2"]=dataframe[["Valeur PNP S2","Valeur PNP 0"]].apply(lambda dataframe: risk(dataframe["Valeur PNP S2"],
                                                                                                                                dataframe["Valeur PNP 0"]),axis=1)
                        dataframe["risque sup CEX S2"]=dataframe[["Valeur CEX S2","Valeur CEX 0"]].apply(lambda dataframe: risk(dataframe["Valeur CEX S2"],
                                                                                                                                dataframe["Valeur CEX 0"]),axis=1)

                        dataframe["risque sup TCH S2"]=dataframe[["Valeur TCH S2","Valeur TCH 0"]].apply(lambda dataframe: risk(dataframe["Valeur TCH S2"],
                                                                                                                                dataframe["Valeur TCH 0"]),axis=1)

                        # Calcul des risques  ponderes
                        dataframe["RWA 0"]=10063541
                        dataframe["RWA Pr√©dicte S2"]=dataframe["risque sup PNP S2"]*0.85+dataframe["risque sup CEX S2"]*0.12+dataframe["risque sup TCH S2"]*0.03+ dataframe["RWA 0"]
                        # ratio de solvabilit√©
                        # st.write("Sc√©nario 3",dataframe)




                        ###"""
                        #### METHODE DE CALCUL SCENARIO 1

                        ####

                        df_s1=pd.DataFrame(list_predict_value_s1,index=None,columns=["PNP S1"])

                        dataframe["PNP S1"]=df_s1["PNP S1"]
                        dfm2_s1=pd.DataFrame(list_predict_m2_s1,index=None, columns=["CEX S1"])
                        dataframe["CEX S1"]=dfm2_s1["CEX S1"]
                        dfm3_s1=pd.DataFrame(list_predict_m3_s1,index=None,columns=["TCH S1"])
                        dataframe["TCH S1"]=dfm3_s1["TCH S1"]

                        dataframe["Variation PNP S1"]=dataframe[["PNP S1","PNP 0"]].apply(lambda dataframe:variation(dataframe["PNP S1"],dataframe["PNP 0"]),axis=1)
                        dataframe["Variation CEX S1"]=dataframe[["CEX S1","CEX 0"]].apply(lambda dataframe:variation(dataframe["CEX S1"],dataframe["CEX 0"]),axis=1)
                        dataframe["Variation TCH S1"]=dataframe[["TCH S1","TCH 0"]].apply(lambda dataframe:variation(dataframe["TCH S1"],dataframe["TCH 0"]),axis=1)

                        # prediction en valeur/ value predict

                        dataframe["Valeur PNP S1"]=dataframe[["Valeur PNP 0","Variation PNP S1"]].apply(lambda dataframe:actualisate(dataframe["Valeur PNP 0"],
                                                                                                                                     dataframe["Variation PNP S1"]),axis=1)
                        dataframe["Valeur CEX S1"]=dataframe[["Valeur CEX 0","Variation CEX S1"]].apply(lambda dataframe:actualisate(dataframe["Valeur CEX 0"],
                                                                                                                                     dataframe["Variation CEX S1"]),axis=1)
                        dataframe["Valeur TCH S1"]=dataframe[["Valeur TCH 0","Variation TCH S1"]].apply(lambda dataframe:actualisate(dataframe["Valeur TCH 0"],
                                                                                                                                     dataframe["Variation TCH S1"]),axis=1)

                        # risque supplementaire/ sup risk

                        dataframe["risque sup PNP S1"]=dataframe[["Valeur PNP S1","Valeur PNP 0"]].apply(lambda dataframe: risk(dataframe["Valeur PNP S1"],
                                                                                                                                dataframe["Valeur PNP 0"]),axis=1)
                        dataframe["risque sup CEX S1"]=dataframe[["Valeur CEX S1","Valeur CEX 0"]].apply(lambda dataframe: risk(dataframe["Valeur CEX S1"],
                                                                                                                                dataframe["Valeur CEX 0"]),axis=1)

                        dataframe["risque sup TCH S1"]=dataframe[["Valeur TCH S1","Valeur TCH 0"]].apply(lambda dataframe: risk(dataframe["Valeur TCH S1"],
                                                                                                                                dataframe["Valeur TCH 0"]),axis=1)

                        # Calcul des risques  ponderes
                        dataframe["RWA 0"]=10063541
                        dataframe["RWA Pr√©dicte S1"]=dataframe["risque sup PNP S1"]*0.85+dataframe["risque sup CEX S1"]*0.12+dataframe["risque sup TCH S1"]*0.03+ dataframe["RWA 0"]
                        # ratio de solvabilit√©
                        # st.write("Sc√©nario 1.2.3",dataframe)

                        # lien de telechargment
                        def generate_excel_download_link_s1(df):
                            towrite = BytesIO()
                            df.to_excel(towrite, encoding="utf-8", index=False, header=True)  # write to BytesIO buffer
                            towrite.seek(0)  # reset pointer
                            b64 = base64.b64encode(towrite.read()).decode()
                            href = f'<a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{b64}" download="Sc√©nario 1.xlsx">Download Excel File</a>'
                            return st.markdown(href, unsafe_allow_html=True)
                        def generate_excel_download_link_s2(df):
                            towrite = BytesIO()
                            df.to_excel(towrite, encoding="utf-8", index=False, header=True)  # write to BytesIO buffer
                            towrite.seek(0)  # reset pointer
                            b64 = base64.b64encode(towrite.read()).decode()
                            href = f'<a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{b64}" download="Sc√©nario 1.xlsx">Download Excel File</a>'
                            return st.markdown(href, unsafe_allow_html=True)

                        def generate_excel_download_link_s3(df):
                            towrite = BytesIO()
                            df.to_excel(towrite, encoding="utf-8", index=False, header=True)  # write to BytesIO buffer
                            towrite.seek(0)  # reset pointer
                            b64 = base64.b64encode(towrite.read()).decode()
                            href = f'<a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{b64}" download="Sc√©nario 1.xlsx">Download Excel File</a>'
                            return st.markdown(href, unsafe_allow_html=True)

                        with sce_1:
                            data_enter=st.number_input("Fonds propres",1000000.00,500000000000.00)

                        st.markdown("***")
                        #with sce_2:
                        # test_choice=st.selectbox("",options=["Select one","Analyse de sc√©nario",
                        #                                               "Analyse de sensibilit√©","NFSR Ratio de liquidit√©"])
                        analyse_a_faire= option_menu(
                            menu_title=None,
                            options=["Analyse de sc√©nario","Analyse de sensibilit√©","NFSR Ratio de liquidit√©"],
                            icons=["reception-4","snow3","speedometer"],
                            default_index=0,
                            orientation="horizontal",
                            styles={
                                "container": {"padding": "0!important","background-color": "#2d2d30"},
                                "icon": {"color": "green", "font-size": "25px"},
                                "nav-link": {"font-size": "20px", "text-align": "left", "margin":"0px"},
                                "nav-link-selected": {"background-color": "#008080"},
                            }
                        )
                        # analyse_a_faire

                        # st.markdown("***")
                        if analyse_a_faire=="Analyse de sc√©nario":
                            dataframe_scenario=dataframe
                            dataframe_scenario["Fonds propres"]=data_enter
                            data_scenario=dataframe_scenario[["Ann√©e","RWA Pr√©dicte","Fonds propres"]]
                            data_scenario["Ratio de solvabilit√©"]= (dataframe_scenario["Fonds propres"]/data_scenario["RWA Pr√©dicte"])*100
                            # st.title("Analyse de sc√©nario")
                            data_scenario_s2=dataframe[["Ann√©e","RWA Pr√©dicte S2","Fonds propres"]]
                            data_scenario_s2.rename(columns={"RWA Pr√©dicte S2":"RWA Pr√©dictes"},inplace=True)
                            data_scenario_s2["Ratio de solvabilit√©"]=(data_scenario_s2["Fonds propres"]/data_scenario_s2["RWA Pr√©dictes"])*100
                            data_scenario_s1=dataframe[["Ann√©e","RWA Pr√©dicte S1","Fonds propres"]]
                            data_scenario_s1.rename(columns={"RWA Pr√©dicte S1":"RWA Pr√©dictes"},inplace=True)
                            data_scenario_s1["Ratio de solvabilit√©"]=(data_scenario_s1["Fonds propres"]/data_scenario_s1["RWA Pr√©dictes"])*100

                            # st.markdown("***")
                            options_expander = st.expander("Je visualise le sc√©nario 1")
                            with options_expander:
                                st.header("Sc√©nario 1")
                                st.table(data_scenario_s1)
                                generate_excel_download_link_s1(data_scenario_s1)


                            options_expander = st.expander("Je visualise le sc√©nario 2")
                            with options_expander:
                                st.header("Sc√©nario 2")
                                st.table(data_scenario_s2)
                                generate_excel_download_link_s2(data_scenario_s2)

                            options_expander = st.expander("Je visualise le sc√©nario 3")
                            with options_expander:
                                st.header("Sc√©nario 3")
                                st.table(data_scenario)
                                generate_excel_download_link_s3(data_scenario)

                        elif analyse_a_faire=="Analyse de sensibilit√©":
                            st.text("Analyse de sensibilit√© est en cours de d√©veloppement")


                        elif analyse_a_faire=="NFSR Ratio de liquidit√©":
                            st.write("NFSR Ratio de liquidit√© est en cours de d√©veloppement")


                            ### --------------------IFRS 9 ---------------------------------

                    elif test_a_faire=="IFRS 9":
                        st.cache(allow_output_mutation=True,show_spinner=True,suppress_st_warning=True)
                        with open('style.css') as f:
                            st.markdown(f'<style>{f.read()}</style>', unsafe_allow_html=True)
                            st.markdown("---")
                        options_expander = st.sidebar.expander("J'importe les donn√©es")
                        with options_expander:
                            st.header("Nom de la banque")
                            nom_utilisateur=st.text_input("")
                            date=st.date_input("Date")

                        choice= option_menu(
                            menu_title=None,
                            options=["PD","LGD","EAD","ECL"],
                            icons=["reception-4","snow3","speedometer","list-task"],
                            menu_icon="cast",
                            default_index=0,
                            orientation="horizontal",


                            styles={
                                "container": {"padding": "0!important"},
                                "icon": {"color": "green", "font-size": "25px"},
                                "nav-link": {"font-size": "14px", "text-align": "left", "margin":"0px"},
                                "nav-link-selected": {"background-color": "#008080"},
                            }
                        )
                        global  df_lgd, df_lgd_s1_ens, df_pd, df_ead
                        with options_expander:
                            data_frame_import = st.file_uploader("T√©l√©charger xlsx",type=["xlsx"])
                        if data_frame_import is not None:
                            with options_expander:
                                st.sidebar.success('T√©l√©chargement reussi!', icon="‚úÖ")
                            try:
                                df_all=pd.read_excel(data_frame_import,sheet_name="LGD")
                                df_pd=pd.read_excel(data_frame_import,sheet_name="PD")
                                # df_ead=pd.read_excel(data_frame_import,sheet_name="EAD")


                                # ------------Marque de saisir ---Enter values--------------------
                                a,s,b=st.sidebar.columns([8,0.2,8])
                                with a:
                                    value_one=st.number_input("Valeur Minimale",90)
                                    value_two=st.number_input("Valeur Maximale",180)

                                options_expander = st.sidebar.expander("Je prends des notes")
                                with options_expander:
                                    st.markdown("""
                                        #### Notes :
                                        * Les valeurs saisies d√©signent les bornes pour la classification des stages.
                                        * 90 et 180 sont les valeurs par d√©faut.
                                        * S1 est inf√©rieur ou √©gal √† la valeur minimale.
                                        * S2 est sup√©rieur √† la valeur minimale et inf√©rieure √† la valeur maximale
                                        * S3 est strictement sup√©rieur √† la valeur maximale                              
                                        """)
                                def bucket1(bucket):
                                    if bucket <= value_one:
                                        return 1
                                    elif bucket > value_one and bucket <= value_two:
                                        return 2
                                    else:
                                        return 3
                                df_all["Bucket"]=df_all["Nombre Jours Impay√©s"].apply(bucket1)
                                df_all["Bucket"]=df_all["Nombre Jours Impay√©s"].apply(bucket1)
                                df_lgd=df_all[["Bucket","Cat√©gorie","Date d'entr√©e","Date de defaut","Encours"]]
                                df_ead=df_all[["Bucket","Cat√©gorie","Montant de pr√™t","Montant re√ßu","Exposition Bilan","Exposition HB"]]
                                with b:
                                    choix_radio_stage=st.radio("",["Strate 1","Strate 2","Strate 3"],
                                                               horizontal=False,index=0)


                            except Exception as e:
                                st.write("T√©l√©charger les donn√©es")
                        try:
                            # --------------PD----------
                            def rating_classe(Nombre):
                                if Nombre <=30:
                                    return "AAA"
                                elif Nombre>30 and Nombre<= 60:
                                    return "AA"
                                elif Nombre >60 and Nombre <= 90:
                                    return "A"
                                elif Nombre >90 and Nombre <= 120:
                                    return "BB"
                                elif Nombre > 120 and Nombre <=150:
                                    return "B"
                                elif Nombre >150 and Nombre <=180:
                                    return "CCC"
                                else:
                                    return "DDD"
                            def fois_100(valeur):
                                return valeur*100

                            # converitr en excel
                            def convert_df(df):
                                return df.to_csv().encode('utf-8')
                            ##

                            def matrice(df):
                                df_pd=df.copy()
                                df_pd["Notation Initiale"]=df_pd["NBJ1"].apply(rating_classe)
                                df_pd["Notation Suivante"]=df_pd["NBJ2"].apply(rating_classe)
                                # matrice de transition
                                transition_matrix=pd.crosstab(df_pd["Notation Initiale"], df_pd["Notation Suivante"])
                                transition_matrix= transition_matrix.div(transition_matrix.sum(axis = 1), axis = 0)
                                transition_matrix=transition_matrix.apply(lambda x: round(x,4))
                                # transition_matrix=transition_matrix.apply(lambda x: fois_100(x))
                                return transition_matrix

                            # ----------Par cat√©gorie ------------------------
                            df_pd_ens=df_pd.copy()
                            df_pd_partucilier=df_pd[df_pd["Cat√©gorie"]=="Particulier"]
                            df_pd_Grande_Entreprise=df_pd[df_pd["Cat√©gorie"]=="Grande Entreprise"]
                            df_pd_pme=df_pd[df_pd["Cat√©gorie"]=="PME"]
                            df_pd_salarie=df_pd[df_pd["Cat√©gorie"]=="Salari√©"]
                            # -------Appliquer la fonction -----

                            df_pd_ens=matrice(df_pd_ens).copy()
                            df_pd_partucilier=matrice(df_pd_partucilier).copy()
                            df_pd_Grande_Entreprise=matrice(df_pd_Grande_Entreprise).copy()
                            df_pd_pme=matrice(df_pd_pme).copy()
                            df_pd_salarie=matrice(df_pd_salarie).copy()
                        #-------------------------Produit matriciel  si n√©cessaire-----------
                            index, colonne=df_pd_ens.index,df_pd_ens.columns
                            # les produits ---------------------*
                            df_pd_ens=np.dot(df_pd_ens.values,df_pd_ens.values)
                            # df_pd_partucilier=np.dot(df_pd_partucilier.values,df_pd_partucilier.values)
                            # df_pd_Grande_Entreprise=np.dot(df_pd_Grande_Entreprise.values, df_pd_Grande_Entreprise.values)
                            # df_pd_pme=np.dot(df_pd_pme.values,df_pd_pme.values)
                            # df_pd_salarie=np.dot(df_pd_salarie.values, df_pd_salarie.values)

                            def matpro(df):
                                df=df.copy()
                                df=pd.DataFrame(data=df, index=index, columns=colonne)
                                df=df.apply(lambda x: fois_100(x))
                                return df
                            df_pd_ens=matpro(df_pd_ens)
                            # df_pd_partucilier=matpro(df_pd_partucilier)
                            # df_pd_Grande_Entreprise=matpro(df_pd_Grande_Entreprise)
                            # df_pd_pme=matpro(df_pd_pme)
                            # df_pd_salarie=matpro(df_pd_salarie)


                    # ----------------------LGD ----------------------------------------------
                            def difference(value_one,value_two):
                                value=int(value_one-value_two)
                                if value >= 0:
                                    return value
                                elif value < 0:
                                    return -value

                            # filtrer
                            df_lgd_s=df_lgd.copy()
                            df_lgd_s1=df_lgd_s[df_lgd_s["Bucket"]==1]
                            df_lgd_s2=df_lgd_s[df_lgd_s["Bucket"]==2]
                            df_lgd_s3=df_lgd_s[df_lgd_s["Bucket"]==3]

                            #--------------------------LGD-------------------------
                            #--------------------stage 1 -----------------------------
                            df_lgd_s1_ens=df_lgd_s1.copy()
                            df_lgd_s1_partucilier=df_lgd_s1[df_lgd_s1["Cat√©gorie"]=="Particulier"]
                            df_lgd_s1_grande_entreprise=df_lgd_s1[df_lgd_s1["Cat√©gorie"]=="Grande Entreprise"]
                            df_lgd_s1_pme=df_lgd_s1[df_lgd_s1["Cat√©gorie"]=="PME"]
                            df_lgd_s1_salarie=df_lgd_s1[df_lgd_s1["Cat√©gorie"]=="Salari√©"]

                            #----------------------------Stage 2-------------------------
                            df_lgd_s2_ens=df_lgd_s2.copy()
                            df_lgd_s2_partucilier=df_lgd_s2[df_lgd_s2["Cat√©gorie"]=="Particulier"]
                            df_lgd_s2_grande_entreprise=df_lgd_s2[df_lgd_s2["Cat√©gorie"]=="Grande Entreprise"]
                            df_lgd_s2_pme=df_lgd_s2[df_lgd_s2["Cat√©gorie"]=="PME"]
                            df_lgd_s2_salarie=df_lgd_s2[df_lgd_s2["Cat√©gorie"]=="Salari√©"]

                            #----------------------------Stage 3-------------------------
                            df_lgd_s3_ens=df_lgd_s3.copy()
                            df_lgd_s3_partucilier=df_lgd_s3[df_lgd_s3["Cat√©gorie"]=="Particulier"]
                            df_lgd_s3_grande_entreprise=df_lgd_s3[df_lgd_s3["Cat√©gorie"]=="Grande Entreprise"]
                            df_lgd_s3_pme=df_lgd_s3[df_lgd_s3["Cat√©gorie"]=="PME"]
                            df_lgd_s3_salarie=df_lgd_s3[df_lgd_s3["Cat√©gorie"]=="Salari√©"]



                            ## ----------------------------------------- Calcul de la LGD--------------------------------------------------------------

                            def calcul_lgd(data):
                                global final_output
                                df_lgd=data

                                df_lgd["colonne"]=df_lgd[["Date d'entr√©e","Date de defaut"]].apply(lambda df_lgd:difference(df_lgd["Date d'entr√©e"],
                                                                                                                            df_lgd["Date de defaut"]),axis=1)
                                triangle=pd.pivot_table(df_lgd, index="Date d'entr√©e",aggfunc="sum",columns="colonne",values="Encours")

                                #--------------------------------------------------------------------------------------------------------------
                                # ----------------------ETAPE DE CALCUL LGD---------------------------------
                                # donn√©es
                                df_lgd=triangle.copy()
                                data_copy=df_lgd.copy()
                                columns_number=len(data_copy.columns)
                                data_copy_lgd=df_lgd.copy()
                                df_soustration=data_copy_lgd
                                #  Soustration faire la difference
                                for i in range(len(df_soustration.columns)-1):
                                    df_soustration["diff{}-{}".format(i+1, i+2)] = df_soustration.iloc[:, i]-df_soustration.iloc[:, i+1]
                                # recuperer la premiere colonne
                                df_soustration.columns._values[0]="Enter"
                                premiere_colonne=df_soustration.iloc[:,0]
                                # recuperer les derni√®res colonnes
                                derniere_colonne=df_soustration.iloc[:,-(len(data_copy.columns)-1):]
                                # # calculer le cumul par ligne de chaque des colonnes
                                df_copy_0=derniere_colonne.copy()
                                derniere_colonne_cumul = pd.DataFrame()
                                n = len(df_copy_0.columns)
                                for i in range(1, n+1):
                                    derniere_colonne_cumul[f'cumul_{i}'] = df_copy_0.iloc[:, :i].sum(axis=1)
                                # r√©cuperer le trianlge de Chain Ladder
                                l = list(derniere_colonne_cumul.keys())
                                for idx, key in enumerate(l[::-1]):
                                    derniere_colonne_cumul[key] = derniere_colonne_cumul[key].iloc[:idx+1]

                                # # combiner les data
                                combiner=pd.concat([premiere_colonne,derniere_colonne_cumul],axis=1)
                                # combiner=pd.concat([premiere_colonne,derniere_colonne],axis=1)
                                # remplacer les valeurs n√©gatives par des zeros
                                for col in combiner.columns:
                                    combiner[col] = combiner[col].apply(lambda x: 0 if x<0 else x)
                                # calculer les rapports.
                                for i in range(len(combiner.columns)-1):
                                    combiner['T {}'.format(i+1)] = combiner.iloc[:, i+1] / combiner.iloc[:, 0]
                                # r√©cup√©rer les rapports
                                rapport=combiner.iloc[:,-(len(data_copy.columns)-1):]
                                # la moyenne
                                moy_rapport= rapport.mean(axis=0)
                                # le taux de recuperation
                                taux_recuperation=pd.DataFrame(moy_rapport,columns=["Taux de r√©cup√©ration"])
                                # taux_recuperation["Taux de recup√©ration"] = taux_recuperation["Valeur"].cumsum()
                                taux_recuperation["Taux LGD"]=1-taux_recuperation["Taux de r√©cup√©ration"]
                                # taux_recuperation["Taux LGD cumul√©"]=taux_recuperation["Taux LGD"].cumsum()
                                # Calculez la variation de la colonne taux de recuperation
                                variation = taux_recuperation["Taux de r√©cup√©ration"].diff()
                                # Calculer le taux de croissance
                                taux_croissance = variation / taux_recuperation["Taux de r√©cup√©ration"].shift(1)
                                taux_croissance=pd.DataFrame(taux_croissance)
                                taux_croissance.columns = ['Taux de croissance']
                                taux_recuperation["Taux de croissance"]=taux_croissance['Taux de croissance']
                                final_output=taux_recuperation.copy()
                                return final_output


                            #---------------------------APPLIQUER LA FONCTION DE LGD  AUX CATEGOIRES----------------------------------------
                            # -------------------------Stage 1----------------------------------------------
                            df_lgd_s1_ens=calcul_lgd(df_lgd_s1_ens).copy()
                            df_lgd_s1_particulier=calcul_lgd(df_lgd_s1_partucilier).copy()
                            df_lgd_s1_grande_entreprise=calcul_lgd(df_lgd_s1_grande_entreprise).copy()
                            df_lgd_s1_pme=calcul_lgd(df_lgd_s1_pme).copy()
                            df_lgd_s1_salarie=calcul_lgd(df_lgd_s1_salarie).copy()


                            # -------------------------Stage 2----------------------------------------------
                            df_lgd_s2_ens=calcul_lgd(df_lgd_s2_ens).copy()
                            df_lgd_s2_particulier=calcul_lgd(df_lgd_s2_partucilier).copy()
                            df_lgd_s2_grande_entreprise=calcul_lgd(df_lgd_s2_grande_entreprise).copy()
                            df_lgd_s2_pme=calcul_lgd(df_lgd_s2_pme).copy()
                            df_lgd_s2_salarie=calcul_lgd(df_lgd_s2_salarie).copy()

                            # -------------------------Stage 3----------------------------------------------
                            df_lgd_s3_ens=calcul_lgd(df_lgd_s3_ens).copy()
                            df_lgd_s3_particulier=calcul_lgd(df_lgd_s3_partucilier).copy()
                            df_lgd_s3_grande_entreprise=calcul_lgd(df_lgd_s3_grande_entreprise).copy()
                            df_lgd_s3_pme=calcul_lgd(df_lgd_s3_pme).copy()
                            df_lgd_s3_salarie=calcul_lgd(df_lgd_s3_salarie).copy()

                    ## -----------------------EAD -------------------------

                            df_ead_s=df_ead.copy()
                            df_ead_s1=df_ead_s[df_ead_s["Bucket"]==1]
                            df_ead_s2=df_ead_s[df_ead_s["Bucket"]==2]
                            df_ead_s3=df_ead_s[df_ead_s["Bucket"]==3]

                            #---------------------------EAD-------------------------------------
                                                        # -------------------Stage 1 --------------------------------------
                            df_ead_s1_ens=df_ead_s1.copy()
                            df_ead_s1_particulier=df_ead_s1[df_ead_s1["Cat√©gorie"]=="Particulier"]
                            df_ead_s1_grande_entreprise=df_ead_s1[df_ead_s1["Cat√©gorie"]=="Grande Entreprise"]
                            df_ead_s1_pme=df_ead_s1[df_ead_s1["Cat√©gorie"]=="PME"]
                            df_ead_s1_salarie=df_ead_s1[df_ead_s1["Cat√©gorie"]=="Salari√©"]

                                                         # -------------------Stage 2 --------------------------------------
                            df_ead_s2_ens=df_ead_s2.copy()
                            df_ead_s2_particulier=df_ead_s2[df_ead_s2["Cat√©gorie"]=="Particulier"]
                            df_ead_s2_grande_entreprise=df_ead_s2[df_ead_s2["Cat√©gorie"]=="Grande Entreprise"]
                            df_ead_s2_pme=df_ead_s2[df_ead_s2["Cat√©gorie"]=="PME"]
                            df_ead_s2_salarie=df_ead_s2[df_ead_s2["Cat√©gorie"]=="Salari√©"]

                                                         # -------------------Stage 3 --------------------------------------
                            df_ead_s3_ens=df_ead_s3.copy()
                            df_ead_s3_particulier=df_ead_s3[df_ead_s3["Cat√©gorie"]=="Particulier"]
                            df_ead_s3_grande_entreprise=df_ead_s3[df_ead_s3["Cat√©gorie"]=="Grande Entreprise"]
                            df_ead_s3_pme=df_ead_s3[df_ead_s3["Cat√©gorie"]=="PME"]
                            df_ead_s3_salarie=df_ead_s3[df_ead_s3["Cat√©gorie"]=="Salari√©"]

                            # ---------------Fonction ----
                            def calcul_ead(df):
                                df=df.copy()
                                # calculer le ccf
                                ccf=(df["Montant de pr√™t"]-df["Montant re√ßu"])/df["Montant de pr√™t"]
                                ccf=ccf.mean()
                                montant=df["Montant de pr√™t"]
                                montant=montant.mean()
                                ead=montant*ccf
                                ead=pd.DataFrame({"Montant":[montant],"CCF":[ccf],"EAD":[ead]})

                                # ead=df["Exposition Bilan"]+ccf*df["Exposition HB"]
                                return ead




                    #---------------------------APPLIQUER LA FONCTION DE EAD  AUX CATEGOIRES----------------------------------------

                                            #----------------------------Stage 1---------------------------------
                            df_ead_s1_ens=calcul_ead(df_ead_s1_ens).copy()
                            df_ead_s1_particulier=calcul_ead(df_ead_s1_particulier).copy()
                            df_ead_s1_grande_entreprise=calcul_ead(df_ead_s1_grande_entreprise).copy()
                            df_ead_s1_pme=calcul_ead(df_ead_s1_pme).copy()
                            df_ead_s1_salarie=calcul_ead(df_ead_s1_salarie).copy()

                                                #----------------------------Stage 2---------------------------------
                            df_ead_s2_ens=calcul_ead(df_ead_s2_ens).copy()
                            df_ead_s2_particulier=calcul_ead(df_ead_s2_particulier).copy()
                            df_ead_s2_grande_entreprise=calcul_ead(df_ead_s2_grande_entreprise).copy()
                            df_ead_s2_pme=calcul_ead(df_ead_s2_pme).copy()
                            df_ead_s2_salarie=calcul_ead(df_ead_s2_salarie).copy()

                                                #-------------------------  Stage 3-------------------------------
                            df_ead_s3_ens=calcul_ead(df_ead_s3_ens).copy()
                            df_ead_s3_particulier=calcul_ead(df_ead_s3_particulier).copy()
                            df_ead_s3_grande_entreprise=calcul_ead(df_ead_s3_grande_entreprise).copy()
                            df_ead_s3_pme=calcul_ead(df_ead_s3_pme).copy()
                            df_ead_s3_salarie=calcul_ead(df_ead_s3_salarie).copy()


                            # --------------------les EAD -----------------
                            index=["Ensemble","Particulier","Grande Entreprise","PME","Salari√©"]
                            df_ead_s1=pd.concat([df_ead_s1_ens,df_ead_s1_particulier,df_ead_s1_grande_entreprise,df_ead_s1_pme,df_ead_s1_salarie],axis=0)
                            df_ead_s2=pd.concat([df_ead_s2_ens,df_ead_s2_particulier,df_ead_s2_grande_entreprise,df_ead_s2_pme,df_ead_s2_salarie],axis=0)
                            df_ead_s3=pd.concat([df_ead_s3_ens,df_ead_s3_particulier,df_ead_s3_grande_entreprise,df_ead_s3_pme,df_ead_s3_salarie],axis=0)

                            df_ead_s1["index"]=index
                            df_ead_s2["index"]=index
                            df_ead_s3["index"]=index
                            df_ead_s1=df_ead_s1.set_index("index")
                            df_ead_s2=df_ead_s2.set_index("index")
                            df_ead_s3=df_ead_s3.set_index("index")


                            if choice=="PD":
                                col1,sep,col2,sep,col3=st.columns([2,0.2,8,0.2,2])
                                with col2:

                                    st.markdown("""
                                        ####  Probabilit√© de d√©faut
                                        """)
                                    choix_radio=st.radio("",["Ensemble","Particulier","Grande Entreprise","PME","Salari√©"],
                                                             horizontal=True,index=0)
                                st.markdown("***")
                                if choix_radio=="Ensemble":
                                    st.table(df_pd_ens)
                                    st.download_button(label="Je t√©l√©charge le fichier", data=convert_df(df_pd_ens),file_name='PD-ENS.csv',mime='text/csv')
                                    #st.write(df_pd_ens.iloc[:,-1].iloc[6])
                                elif choix_radio=="Particulier":
                                    st.table(df_pd_partucilier)
                                    st.download_button(label="Je t√©l√©charge le fichier", data=convert_df(df_pd_partucilier),file_name='PD-PARTICULIER.csv',mime='text/csv')
                                elif choix_radio=="Grande Entreprise":
                                    st.table(df_pd_Grande_Entreprise)
                                    st.download_button(label="Je t√©l√©charge le fichier", data=convert_df(df_pd_Grande_Entreprise),file_name='PD-GEE.csv',mime='text/csv')

                                elif choix_radio=="PME":
                                    st.table(df_pd_pme)
                                    st.download_button(label="Je t√©l√©charge le fichier", data=convert_df(df_pd_pme),file_name='PD-PME.csv',mime='text/csv')
                                elif choix_radio=="Salari√©":
                                    st.table(df_pd_salarie)
                                    st.download_button(label="Je t√©l√©charge le fichier", data=convert_df(df_pd_salarie),file_name='PD-SALARIE.csv',mime='text/csv')

                            # afficher les r√©sultats de LGD

                            elif choice=="LGD":
                                col1,sep,col2,sep,col3=st.columns([2,0.2,8,0.2,2])
                                with col2:

                                    st.markdown("""
                                        ####  Perte en cas de d√©faut : Loss Given Default
                                        """)
                                    choix_radio=st.radio("",["Ensemble","Particulier","Grande Entreprise","PME","Salari√©"],
                                                             horizontal=True,index=0)

                                st.markdown("***")

                                # ----------------Stage 1----------------------------------
                                if choix_radio_stage=="Strate 1":
                                    if choix_radio=="Ensemble":
                                        st.table(df_lgd_s1_ens)
                                        st.download_button(label="Je t√©l√©charge le fichier", data=convert_df(df_lgd_s1_ens),file_name='LGD-ENS-S1.csv',mime='text/csv')
                                        #st.write(df_ead_s1.iloc[:,-1].iloc[0])
                                    elif choix_radio=="Particulier":
                                        st.table(df_lgd_s1_particulier)
                                        st.download_button(label="Je t√©l√©charge le fichier", data=convert_df(df_lgd_s1_particulier),file_name='LGD-PARTICULIER-S1.csv',mime='text/csv')
                                    elif choix_radio=="Grande Entreprise":
                                        st.table(df_lgd_s1_grande_entreprise)
                                        st.download_button(label="Je t√©l√©charge le fichier", data=convert_df(df_lgd_s1_grande_entreprise),file_name='LGD-GE-S1.csv',mime='text/csv')
                                    elif choix_radio=="PME":
                                        st.table(df_lgd_s1_pme)
                                        st.download_button(label="Je t√©l√©charge le fichier", data=convert_df(df_lgd_s1_pme),file_name='LGD-PME-S1.csv',mime='text/csv')

                                    elif choix_radio=="Salari√©":
                                        st.table(df_lgd_s1_salarie)
                                        st.download_button(label="Je t√©l√©charge le fichier", data=convert_df(df_lgd_s1_salarie),file_name='LGD-SALARIE-S1.csv',mime='text/csv')

                                    else:
                                        st.write("Veuillez choisir une option")

                                #--------------------Stage 2---------------------------------------------
                                elif choix_radio_stage=="Strate 2":
                                    if choix_radio=="Ensemble":
                                        st.table(df_lgd_s2_ens)
                                        st.download_button(label="Je t√©l√©charge le fichier", data=convert_df(df_lgd_s2_ens),file_name='LGD-ENS-S2.csv',mime='text/csv')
                                    elif choix_radio=="Particulier":
                                        st.table(df_lgd_s2_particulier)
                                        st.download_button(label="Je t√©l√©charge le fichier", data=convert_df(df_lgd_s2_particulier),file_name='LGD-Particulier-S2.csv',mime='text/csv')
                                    elif choix_radio=="Grande Entreprise":
                                        st.table(df_lgd_s2_grande_entreprise)
                                        st.download_button(label="Je t√©l√©charge le fichier", data=convert_df(df_lgd_s2_grande_entreprise),file_name='LGD-GE-S2.csv',mime='text/csv')
                                    elif choix_radio=="PME":
                                        st.table(df_lgd_s2_pme)
                                        st.download_button(label="Je t√©l√©charge le fichier", data=convert_df(df_lgd_s2_pme),file_name='LGD-PME-S2.csv',mime='text/csv')

                                    elif choix_radio=="Salari√©":
                                        st.table(df_lgd_s2_salarie)
                                        st.download_button(label="Je t√©l√©charge le fichier", data=convert_df(df_lgd_s2_salarie),file_name='LGD-Salari√©-S2.csv',mime='text/csv')
                                    else:
                                        st.write("Veuillez choisir une option")

                                #---------------------Stage 3-----------------------------------------------
                                elif choix_radio_stage=="Strate 3":
                                    if choix_radio=="Ensemble":
                                        st.table(df_lgd_s3_ens)
                                        st.download_button(label="Je t√©l√©charge le fichier", data=convert_df(df_lgd_s3_ens),file_name='LGD-ENS-S3.csv',mime='text/csv')
                                    elif choix_radio=="Particulier":
                                        st.table(df_lgd_s3_particulier)
                                        st.download_button(label="Je t√©l√©charge le fichier", data=convert_df(df_lgd_s3_particulier),file_name='LGD-Particulier-S3.csv',mime='text/csv')
                                    elif choix_radio=="Grande Entreprise":
                                        st.table(df_lgd_s3_grande_entreprise)
                                        st.download_button(label="Je t√©l√©charge le fichier", data=convert_df(df_lgd_s3_grande_entreprise),file_name='LGD-GE-S3.csv',mime='text/csv')
                                    elif choix_radio=="PME":
                                        st.table(df_lgd_s3_pme)
                                        st.download_button(label="Je t√©l√©charge le fichier", data=convert_df(df_lgd_s3_pme),file_name='LGD-PME-S3.csv',mime='text/csv')

                                    elif choix_radio=="Salari√©":
                                        st.table(df_lgd_s3_salarie)
                                        st.download_button(label="Je t√©l√©charge le fichier", data=convert_df(df_lgd_s3_salarie),file_name='LGD-Salari√©-S3.csv',mime='text/csv')

                                    else:
                                        st.write("Veuillez choisir une option")


                            elif choice=="EAD":
                                col1,sep,col2,sep,col3=st.columns([2,0.2,8,0.2,2])
                                with col2:
                                    # st.markdown("***")
                                    st.markdown("""
                                        ####  Exposition en cas de defaut: Expositon at default
                                        """)
                                    # choix_radio_ead=st.radio("",["Ensemble","Particulier","Grande Entreprise","PME","Salari√©"],
                                    #                          horizontal=True,index=0)

                                st.markdown("***")
                                if choix_radio_stage=="Strate 1":
                                    st.header("Stage 1")
                                    st.table(df_ead_s1)
                                    st.download_button(label="Je t√©l√©charge le fichier", data=convert_df(df_ead_s1),file_name='EAD-S1.csv',mime='text/csv')
                                elif choix_radio_stage=="Strate 2":
                                    st.header("Stage 2")
                                    st.table(df_ead_s2)
                                    st.download_button(label="Je t√©l√©charge le fichier", data=convert_df(df_ead_s2),file_name='EAD-S2.csv',mime='text/csv')
                                elif choix_radio_stage=="Strate 3":
                                    st.header("Stage 3")
                                    st.table(df_ead_s3)
                                    st.download_button(label="Je t√©l√©charge le fichier", data=convert_df(df_ead_s3),file_name='EAD-S3.csv',mime='text/csv')

                                else:
                                    st.write("Veuillez choisir une option")

                            elif choice=="ECL":
                                col1,sep,col2,sep,col3=st.columns([2,0.2,8,0.2,2])
                                with col2:

                                    # st.markdown("***")
                                    st.markdown("""
                                            ####  Perte attendues de cr√©dits : Expected credit loss
                                            """)
                                    # choix=st.radio("",["Ensemble","Particulier","Grande Entreprise","PME","Salari√©"],
                                    #                          horizontal=True,index=0)
                                st.markdown("***")

                                if choix_radio_stage=="Strate 1":
                                    st.header("Stage 1")
                                    # st.write((df_ead_s1.iloc[:,-1]).iloc[0])
                                    # calcul de PD * EAD ? ici iloc permet de r√©cup√©rer les valeurs qu'on souhaite et
                                    # qui correspond parfaitement au nom
                                    index_pd=(df_pd_ens.iloc[:,-1]).iloc[0]/100*((df_ead_s1.iloc[:,-1]).iloc[0])
                                    index_pd_p=(df_pd_partucilier.iloc[:,-1]).iloc[0]/100*((df_ead_s1.iloc[:,-1]).iloc[1])
                                    index_pd_g=(df_pd_Grande_Entreprise.iloc[:,-1]).iloc[0]/100*((df_ead_s1.iloc[:,-1]).iloc[2])
                                    index_pd_pme=(df_pd_pme.iloc[:,-1]).iloc[0]/100*((df_ead_s1.iloc[:,-1]).iloc[3])
                                    index_pd_s=(df_pd_salarie.iloc[:,-1]).iloc[0]/100*((df_ead_s1.iloc[:,-1]).iloc[4])

                                    # ici on reinitialise l'index de LGD, et on cr√©e une table LGD

                                    df_ecl=df_lgd_s1_ens["Taux LGD"].reset_index()
                                    df1=df_lgd_s1_particulier["Taux LGD"].reset_index()
                                    df_ecl["LGD P"]=df1["Taux LGD"]
                                    df1=df_lgd_s1_grande_entreprise["Taux LGD"].reset_index()
                                    df_ecl["LGD G"]=df1["Taux LGD"]
                                    df1=df_lgd_s1_pme["Taux LGD"].reset_index()
                                    df_ecl["LGD PM"]=df1["Taux LGD"]
                                    df1=df_lgd_s1_salarie["Taux LGD"].reset_index()
                                    df_ecl["LGD S"]=df1["Taux LGD"]

                                    # Ajouter le produit PD* EAD  √† la table LGD
                                    df_ecl["PD EAD E"]=index_pd
                                    df_ecl["PD EAD P"]=index_pd_p
                                    df_ecl["PD EAD G"]=index_pd_g
                                    df_ecl["PD EAD PM"]=index_pd_pme
                                    df_ecl["PD EAD S"]=index_pd_s
                                    # calcul du ecl
                                    df_ecl["ECL E"]=df_ecl["Taux LGD"]*df_ecl["PD EAD E"]
                                    df_ecl["ECL P"]=df_ecl["LGD P"]*df_ecl["PD EAD P"]
                                    df_ecl["ECL G"]=df_ecl["LGD G"]*df_ecl["PD EAD G"]
                                    df_ecl["ECL PME"]=df_ecl["LGD PM"]*df_ecl["PD EAD PM"]
                                    df_ecl["ECL S"]=df_ecl["LGD S"]*df_ecl["PD EAD S"]
                                    df_ecl=df_ecl[["ECL E","ECL P","ECL G","ECL PME","ECL S","index"]]
                                    df_ecl=df_ecl.set_index("index")
                                    st.table(df_ecl)
                                    telech=convert_df(df_ecl)
                                    st.download_button(label="Je t√©l√©charge le fichier", data=telech,file_name='ECL-S1.csv',mime='text/csv')

                                elif choix_radio_stage=="Strate 2":
                                    st.header("Stage 2")
                                    index_pd=(df_pd_ens.iloc[:,-1]).iloc[5]/100*((df_ead_s2.iloc[:,-1]).iloc[0])
                                    index_pd_p=(df_pd_partucilier.iloc[:,-1]).iloc[5]/100*((df_ead_s2.iloc[:,-1]).iloc[1])
                                    index_pd_g=(df_pd_Grande_Entreprise.iloc[:,-1]).iloc[5]/100*((df_ead_s2.iloc[:,-1]).iloc[2])
                                    index_pd_pme=(df_pd_pme.iloc[:,-1]).iloc[5]/100*((df_ead_s2.iloc[:,-1]).iloc[3])
                                    index_pd_s=(df_pd_salarie.iloc[:,-1]).iloc[5]/100*((df_ead_s2.iloc[:,-1]).iloc[4])

                                    df_ecl=df_lgd_s2_ens["Taux LGD"].reset_index()
                                    df1=df_lgd_s2_particulier["Taux LGD"].reset_index()
                                    df_ecl["LGD P"]=df1["Taux LGD"]
                                    df1=df_lgd_s2_grande_entreprise["Taux LGD"].reset_index()
                                    df_ecl["LGD G"]=df1["Taux LGD"]
                                    df1=df_lgd_s2_pme["Taux LGD"].reset_index()
                                    df_ecl["LGD PM"]=df1["Taux LGD"]
                                    df1=df_lgd_s2_salarie["Taux LGD"].reset_index()
                                    df_ecl["LGD S"]=df1["Taux LGD"]

                                    df_ecl["PD EAD E"]=index_pd
                                    df_ecl["PD EAD P"]=index_pd_p
                                    df_ecl["PD EAD G"]=index_pd_g
                                    df_ecl["PD EAD PM"]=index_pd_pme
                                    df_ecl["PD EAD S"]=index_pd_s
                                    # le ecl
                                    df_ecl["ECL E"]=df_ecl["Taux LGD"]*df_ecl["PD EAD E"]
                                    df_ecl["ECL P"]=df_ecl["LGD P"]*df_ecl["PD EAD P"]
                                    df_ecl["ECL G"]=df_ecl["LGD G"]*df_ecl["PD EAD G"]
                                    df_ecl["ECL PME"]=df_ecl["LGD PM"]*df_ecl["PD EAD PM"]
                                    df_ecl["ECL S"]=df_ecl["LGD S"]*df_ecl["PD EAD S"]
                                    df_ecl=df_ecl[["ECL E","ECL P","ECL G","ECL PME","ECL S","index"]]
                                    df_ecl=df_ecl.set_index("index")
                                    st.table(df_ecl)
                                    telech=convert_df(df_ecl)
                                    st.download_button(label="Je t√©l√©charge le fichier", data=telech,file_name='ECL-S2.csv',mime='text/csv')

                                elif choix_radio_stage=="Strate 3":
                                    st.header("Stage 3")

                                    index_pd=(df_pd_ens.iloc[:,-1]).iloc[6]/100*((df_ead_s3.iloc[:,-1]).iloc[0])
                                    index_pd_p=(df_pd_partucilier.iloc[:,-1]).iloc[6]/100*((df_ead_s3.iloc[:,-1]).iloc[1])
                                    #st.write(df_pd_Grande_Entreprise.iloc[:,-1].iloc[6])
                                    index_pd_g=(df_pd_Grande_Entreprise.iloc[:,-1]).iloc[5]/100*((df_ead_s3.iloc[:,-1]).iloc[2])

                                    #st.write(index_pd)
                                    index_pd_pme=(df_pd_pme.iloc[:,-1]).iloc[6]/100*((df_ead_s3.iloc[:,-1]).iloc[3])
                                    #st.write(index_pd)
                                    index_pd_s=(df_pd_salarie.iloc[:,-1]).iloc[6]/100*((df_ead_s3.iloc[:,-1]).iloc[4])
                                    #st.write(index_pd)

                                    df_ecl=df_lgd_s3_ens["Taux LGD"].reset_index()
                                    df1=df_lgd_s3_particulier["Taux LGD"].reset_index()
                                    df_ecl["LGD P"]=df1["Taux LGD"]
                                    df1=df_lgd_s3_grande_entreprise["Taux LGD"].reset_index()
                                    df_ecl["LGD G"]=df1["Taux LGD"]
                                    df1=df_lgd_s3_pme["Taux LGD"].reset_index()
                                    df_ecl["LGD PM"]=df1["Taux LGD"]
                                    df1=df_lgd_s3_salarie["Taux LGD"].reset_index()
                                    df_ecl["LGD S"]=df1["Taux LGD"]

                                    df_ecl["PD EAD E"]=index_pd
                                    df_ecl["PD EAD P"]=index_pd_p
                                    df_ecl["PD EAD G"]=index_pd_g
                                    df_ecl["PD EAD PM"]=index_pd_pme
                                    df_ecl["PD EAD S"]=index_pd_s
                                    # le ecl
                                    df_ecl["ECL E"]=df_ecl["Taux LGD"]*df_ecl["PD EAD E"]
                                    df_ecl["ECL P"]=df_ecl["LGD P"]*df_ecl["PD EAD P"]
                                    df_ecl["ECL G"]=df_ecl["LGD G"]*df_ecl["PD EAD G"]
                                    df_ecl["ECL PME"]=df_ecl["LGD PM"]*df_ecl["PD EAD PM"]
                                    df_ecl["ECL S"]=df_ecl["LGD S"]*df_ecl["PD EAD S"]
                                    df_ecl=df_ecl[["ECL E","ECL P","ECL G","ECL PME","ECL S","index"]]
                                    df_ecl=df_ecl.set_index("index")
                                    st.table(df_ecl)
                                    #--------------------------------------Rapport pptx-----------------------------------------
                                    # # ----------------------------------Cr√©er une pr√©sentation PowerPoint------------------------
                                    # prs = Presentation()
                                    # # ----------------------------------Ajouter une diapositive avec un titre et un sous-titre
                                    #

                                    telech=convert_df(df_ecl)
                                    st.download_button(label="Je t√©l√©charge le fichier", data=telech,file_name='ECL-S3.csv',mime='text/csv')
                                    # slide = prs.slides.add_slide(prs.slide_layouts[0])
                                    # title = slide.shapes.title
                                    # title.text = f"Mise en place op√©rationnelle de la norme IFRS 9 au sein de la banque {nom_utilisateur}"
                                    # subtitle = slide.placeholders[1]
                                    # subtitle.text = "Rapport de Mission"
                                    # # Ajouter une diapositive avec une image
                                    # img_path = 'Y3c.jpg'
                                    # left = top = Inches(0.1)
                                    # height = Inches(2)
                                    # pic = slide.shapes.add_picture(img_path, left, top, height=height)
                                    #



                            else:
                                st.write("Veuillez choisir une option")


                        except Exception as e:
                            st.write("Veuillez saisir les donn√©es")

##
                    #  ----------REVERSE STRESS TEST ---------------------------------------------------------------------
                    # #

                    elif test_a_faire=="Reverses Tests":
                        st.cache(allow_output_mutation=True,show_spinner=True,suppress_st_warning=True)
                        # Cr√©er un expander dans la barre lat√©rale
                        options_expander = st.sidebar.expander("J'importe les donn√©es")
                        with options_expander:
                            st.subheader("Param√®tre de donn√©es")
                        # ajouter untitre et une image le logo de y3
                        st.title("REVERSE STRESS TEST : RISK MANAGEMENT TOOL ")
                        st.markdown("***")
                        # methode 2
                        a1, sep, a2,sep,a3=st.columns([8,0.2,8,0.2,8])
                        with a1:
                            pnp=st.slider("Valeur de perte", min_value=4.00, max_value=30.00)

                        with a2:
                            credittotal=st.number_input("Cr√©dit total", min_value=100000.00, max_value=1000000000000.00)
                            pret_a_risque=pnp * credittotal/100
                            st.write("Cr√©dit √† risque", pret_a_risque)
                        with a3:
                            fp=st.number_input("Fonds propres", min_value=1000.00)
                            rs=np.round(fp/pret_a_risque*100,2)
                            st.write("Ratio de solvabilit√©",rs,"%")
                        st.markdown("***")
                        L=pnp
                        with options_expander:
                        ### se connecter au donn√©es"""
                            data_file = st.file_uploader("T√©l√©charger CSV ou xlsx",type=["csv","xlsx"])
                        d1,sep,d2,sep,d3=st.columns([8,0.2,8,0.2,8])
                        c1, sep, c2=st.columns([8,0.5,8,])
                        # global df
                        global cov_matrix
                        global scenarios
                        if data_file is not None:
                            with options_expander:
                                st.success("T√©l√©chargement reussi")
                            try:
                                df=pd.read_csv(data_file)
                            except Exception as e:
                                df=pd.read_excel(data_file)


                        try:
                            # st.sidebar.write("Les 5 premi√®res colonnes",df.head(5))
                            with d1:
                                colomns=st.multiselect("Selectionner les variables pour cr√©er les sc√©narios", df.columns)
                                st.markdown("***")
                            df=df[colomns]


                            st.sidebar.write("Donn√©es historiques pour les sc√©narios",df)
                        except Exception as e:
                            st.write("Veuillez t√©l√©charger les donn√©es")


                        try:
                            cov_matrix = pd.DataFrame.cov(df)
                        except Exception as e:
                            st.write("!!!")
                        ### Definition des fonctions"""
                        def gram_schmidt_columns(x):
                            q, r = np.linalg.qr(x)   # caclul de la factorisation de matrice # decomposition
                            return q

                        def gram_schmidt_rows(x):
                            """ row-wise Gram Schmidt procedure """
                            q, r = np.linalg.qr(x.T)
                            return q.T

                        def p_from_w(w):
                            """ √©tant donn√© le vecteur de poids w, calcul√© la matrice de transformation p,
                            comme decrite √† l'annexe  A du livre
                            """
                            w = np.array(w) # table de w
                            n = len(w) # la longueur de w
                            p = np.zeros((n, n)) # une matrice carr√©e nulle de n*n dimension
                            p[0, :] = w / np.linalg.norm(w) # norme de matrice ou vectorielle
                            j0 = (np.abs(w) > 0.00001).nonzero()[0][0]
                            for j in range(j0):
                                p[j + 1, j] = 1
                            for j in range(j0 + 1, n):
                                p[j, j] = 1
                            gs = gram_schmidt_rows(p)
                            for j in range(gs.shape[1]):
                                if gs[0, j] < 0:
                                    gs[:, j] = -gs[:, j]
                            return gs

                        def standard_simplex_vertices(M):
                            """ just that """
                            z = np.eye(M) - 1 / M
                            a = np.eye(M)
                            a[0, :] = 1
                            p = gram_schmidt_rows(a)
                            cls = np.append(range(1, M), 0)
                            gs = np.delete(p.dot(z), 0, axis=0)[:, cls]
                            if gs[0, 0] < 0:
                                gs = -gs
                            return gs * np.sqrt(M / (M - 1))


                        def quadratic_form_inv(a, x):
                            """ x * a^(-1) * x """
                            ax = np.linalg.solve(a, x)
                            return x.T.dot(ax)

                        ### procedure principale"""

                        def comp_scenarios(c, w, L, q, M, printIt=False):
                            """ matrice de covariance  c, vectur poids (param√®tre) w,
                            perte L, vraissemblance relative q,
                            nombre de sc√©nario non-central M,
                            calcul de stress scenarios, et de vraissemblance relative;

                            """
                            # scenario central
                            cw = c.dot(w)
                            ah = L * cw / w.dot(cw)
                            if printIt:
                                print('------------   ah   ------------')
                                print(ah.round(3))

                            # pr√©parer la matrice de transformation
                            p = p_from_w(w)
                            n = len(w)
                            if printIt:
                                print('------------   p   ------------')
                                print(p.round(3))

                            ahy_full = p.dot(ah)
                            first_comp = ahy_full[0]
                            if printIt:
                                print('------------   ahy_full   ------------')
                                print(ahy_full.round(3))

                            # asset names
                            ass = c.columns

                            # compute covariance matrix in the transformed space
                            dy = p.dot(c.dot(p.T))
                            if printIt:
                                print('------------   dy   ------------')
                                print(dy.round(3))

                            # extract decomposition parts
                            d11 = dy[0, 0]
                            # d1I = dy[0, 1:]
                            dI1 = dy[1:, 0]
                            dII = dy[1:, 1:]

                            # compute center and conditional distribution in transformed space
                            wn = np.linalg.norm(w)
                            ahy = L * dI1 / d11 / wn
                            i_m = np.matrix(dI1).reshape(n - 1, 1) * np.matrix(dI1).reshape(1, n - 1)
                            dyh = dII - i_m / d11
                            if printIt:
                                print('------------   ahy   ------------')
                                print(ahy.round(3))
                                print('------------   dyh   ------------')
                                print(dyh.round(3))

                            # egenvalues in ascending order and eigenvectors
                            vals, vect = np.linalg.eig(dyh)
                            ind = vals.argsort()[::-1]
                            V = vect[:, ind]

                            # place the eigenvectors in opposite order
                            # V = vect[:, ::-1]
                            if printIt:
                                print('------------   V   ------------')
                                print(V.round(3))
                                print('------------   Vals   ------------')
                                print(vals.round(3))

                            # compute vertices of the regular M-simplex
                            t = standard_simplex_vertices(M)
                            if printIt:
                                print('------------   t   ------------')
                                print(t.round(3))

                            # rotation to principle components space
                            z = V[:, :(M - 1)].dot(t)
                            if printIt:
                                print('------------   z   ------------')
                                print(z.round(3))

                            # compute the sphere radius
                            z_last = z[:, -1]
                            denom = quadratic_form_inv(dyh, np.array(z_last))
                            r = np.sqrt(-2 * np.log(q) / denom)[0, 0]
                            ahyr = ahy.repeat(M).reshape(n - 1, M)
                            rz = r * z + ahyr
                            if printIt:
                                print('------------   r   ------------')
                                print(r.round(3))

                            # unite central scenario with sphere scenarios
                            df = pd.DataFrame(ahy, columns=['aa']).join(pd.DataFrame(rz))
                            cls = ['Scenario_{0}'.format(i) for i in range(M + 1)]
                            df.columns = cls

                            # append first_comp row at the beginning
                            df0 = first_comp * pd.DataFrame(np.ones((1, M + 1)), columns=cls)
                            df = df0.append(df.set_index(np.array(range(1, n))))
                            if printIt:
                                print('------------   df0   ------------')
                                print(df0.round(3))
                                print('------------   df   ------------')
                                print(df.round(3))

                            # rotating back to initial space
                            ar = p.T.dot(df)

                            # turning to dataframe
                            scenarios = pd.DataFrame(ar, columns=cls).set_index(ass)
                            if printIt:
                                print('------------   scenarios  ------------')
                                print(scenarios.round(3))

                            lik = np.array([multivariate_normal.pdf(scenarios[cls[j]],
                                                                    mean=np.zeros((n,)),
                                                                    cov=c)
                                            for j in range(M + 1)])
                            lik = pd.DataFrame((lik / lik[0]).reshape(1, M + 1), columns=cls)
                            lik.index = ['Likelihood']
                            if printIt:
                                print('------------   lik  ------------')
                                print(lik.round(3))

                            return scenarios, lik

                        try:
                            c =cov_matrix
                            w = np.ones((len(cov_matrix),))
                            M =df.columns.nunique()
                            q = 0.01
                            scenarios, lik = comp_scenarios(c, w, L, q, M)
                            d=df
                            d.tail()
                            #print(lik.round(3))
                            with c2:
                                st.markdown("***")
                                st.header("La probabilit√© de vraissemble")
                                st.table(lik.round(3)) # la vraissemblance relative
                                st.markdown("***")
                            with c1:
                                st.markdown("***")
                                st.header("Sc√©narios G√©n√©r√©s")
                                st.table(scenarios) # sc√©narios

                                def convert_df(df):
                                    return df.to_csv().encode('utf-8')
                                telech=convert_df(scenarios)
                                st.download_button(label="T√©l√©charger ici", data=telech,file_name='scenario.csv',mime='text/csv')
                                #telecharger=st.download_button("T√©l√©charger le r√©sultat",data=telech,file_name="R√©sultats des scenario.xlsx")
                                st.markdown("***")

                            distance=np.array([np.linalg.norm(scenarios.iloc[:, i] - scenarios.iloc[:, j]) for i in range(M) for j in range(i+1, M+1)]).round(3) # la distance

                        except Exception as e:
                            print("Veuillez t√©l√©charger les donn√©es")

                        # no necessary, mais permet d'avoir plus de controle sur l'app
                    else:
                        try:
                            st.write("Veuillez selectionner une option")
                        except Exception as e:
                            st.write("Veuillez selectionner une option")

                else:
                    with st.spinner("Conexion en cours..."):
                        time.sleep(1)
                        st.warning("l'identifiant ou le mot de passe est incorrect")


    elif choice=="S'inscrire":
        st.subheader("Cr√©ation d'un nouveau compte")
        new_user = st.text_input("Nom d'utilisateur")
        new_password = st.text_input("Mot de passe",type='password')
        if st.button("S'inscrire"):
            if new_password and new_user is not None:
                create_usertable()
                add_userdata(new_user,make_hashes(new_password))
                with st.spinner("Traitement en cours..."):
                    time.sleep(2)
                    r1=st.success("Votre compte a √©t√© cr√©e avec succ√®s !")
                    r2=st.info("Vous pouvez vous connecter maintenant")
                    for i in range(2):
                        time.sleep(1)
                    r1.empty()
                    r2.empty()


            else:

                # Afficher une alerte de succ√®s pendant 3 secondes
                with st.spinner("Traitement en cours..."):
                    time.sleep(2)
                    result = st.warning("Veuillez remplir toutes les options")

                    # Attendre 3 secondes avant de masquer l'alerte
                    for i in range(2):
                        time.sleep(1)

                    # Effacer l'alerte de succ√®s
                    result.empty()


if __name__ == '__main__':
    main()